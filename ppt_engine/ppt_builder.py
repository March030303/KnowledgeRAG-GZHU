#!/usr/bin/env python3
"""
PPTBuilder - Core PPT generation engine using python-pptx
Supports multiple layout types and cyberpunk neon styling
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from typing import List, Dict, Optional, Tuple, Any
import os

# ============================================================
# Color Palette - Cyberpunk Neon Dark
# ============================================================
BG_DARK = RGBColor(0x12, 0x12, 0x17)       # #121217
CARD_BG = RGBColor(0x1E, 0x1E, 0x24)       # #1e1e24
TEXT_PRIMARY = RGBColor(0xFF, 0xFF, 0xFF)  # #ffffff
TEXT_SECONDARY = RGBColor(0xE5, 0xE7, 0xEB) # #e5e7eb
TEXT_MUTED = RGBColor(0x9C, 0xA3, 0xAF)    # #9ca3af
NEON_PURPLE = RGBColor(0xA8, 0x55, 0xF7)   # #a855f7
NEON_PINK = RGBColor(0xEC, 0x48, 0x99)     # #ec4899
NEON_EMERALD = RGBColor(0x34, 0xD3, 0x99)  # #34d399
NEON_BLUE = RGBColor(0x60, 0xA5, 0xFA)     # #60a5fa
NEON_ORANGE = RGBColor(0xFB, 0x92, 0x3C)   # #fb923c
NEON_RED = RGBColor(0xEF, 0x44, 0x44)      # #ef4444

# Chapter accent colors
CHAPTER_COLORS = {
    1: NEON_PURPLE,
    2: NEON_EMERALD,
    3: NEON_ORANGE,
    4: NEON_BLUE,
    5: NEON_PINK,
}


class PPTBuilder:
    """PPT构建器 - 封装python-pptx操作"""
    
    def __init__(self, template_path: Optional[str] = None):
        """初始化PPT构建器
        
        Args:
            template_path: 可选的模板文件路径(.potx)
        """
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        # Set 16:9 widescreen
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
        self.slide_count = 0
    
    def _set_slide_bg(self, slide, color: RGBColor = BG_DARK):
        """设置幻灯片背景色"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
    
    def _add_text_box(self, slide, left, top, width, height, 
                      text: str, font_size: int = 18, bold: bool = False,
                      color: RGBColor = TEXT_PRIMARY, 
                      align: PP_ALIGN = PP_ALIGN.LEFT,
                      font_name: str = "Microsoft YaHei"):
        """添加文本框"""
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align
        return shape
    
    def _add_rounded_rect(self, slide, left, top, width, height,
                          fill_color: RGBColor, 
                          line_color: Optional[RGBColor] = None,
                          line_width: Pt = Pt(1)):
        """添加圆角矩形"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
            shape.line.width = line_width
        else:
            shape.line.fill.background()
        return shape
    
    def _add_bullet_text(self, slide, left, top, width, height,
                        items: List[str], font_size: int = 14,
                        color: RGBColor = TEXT_SECONDARY,
                        bullet_color: RGBColor = NEON_PURPLE):
        """添加带项目符号的文本"""
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"▸  {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = "Microsoft YaHei"
            p.space_after = Pt(8)
        return shape
    
    def add_cover(self, title: str, subtitle: str, description: str,
                  accent: RGBColor = NEON_PURPLE):
        """添加封面页"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_bg(slide)
        
        # Title
        self._add_text_box(slide, Inches(0), Inches(2.5), Inches(13.333), Inches(1.5),
                          title, font_size=72, bold=True, 
                          color=TEXT_PRIMARY, align=PP_ALIGN.CENTER)
        
        # Subtitle
        self._add_text_box(slide, Inches(0), Inches(4.2), Inches(13.333), Inches(0.8),
                          subtitle, font_size=32, 
                          color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
        
        # Description
        self._add_text_box(slide, Inches(0), Inches(5.0), Inches(13.333), Inches(0.6),
                          description, font_size=20, 
                          color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        
        # Neon accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                      Inches(5.5), Inches(5.8), 
                                      Inches(2.333), Pt(4))
        line.fill.solid()
        line.fill.fore_color.rgb = accent
        line.line.fill.background()
        
        self.slide_count += 1
        return slide
    
    def add_toc(self, title: str, items: List[Tuple[str, str, str, RGBColor]]):
        """添加目录页
        
        Args:
            items: List of (number, title, description, accent_color)
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_bg(slide)
        
        # Title
        self._add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.8),
                          title, font_size=36, bold=True, color=TEXT_PRIMARY)
        
        # Cards
        card_width = Inches(5.8)
        card_height = Inches(2.2)
        gap_x = Inches(0.4)
        gap_y = Inches(0.4)
        start_x = Inches(0.5)
        start_y = Inches(1.5)
        
        for i, (num, item_title, desc, accent) in enumerate(items):
            row = i // 2
            col = i % 2
            x = start_x + col * (card_width + gap_x)
            y = start_y + row * (card_height + gap_y)
            
            # Card background
            card = self._add_rounded_rect(slide, x, y, card_width, card_height, 
                                         CARD_BG, accent, Pt(1.5))
            
            # Number
            self._add_text_box(slide, x + Inches(0.2), y + Inches(0.15), 
                              Inches(1.5), Inches(0.8),
                              num, font_size=48, bold=True, 
                              color=accent, align=PP_ALIGN.LEFT)
            
            # Title
            self._add_text_box(slide, x + Inches(0.2), y + Inches(0.9), 
                              Inches(5.4), Inches(0.5),
                              item_title, font_size=24, bold=True, 
                              color=TEXT_PRIMARY)
            
            # Description
            self._add_text_box(slide, x + Inches(0.2), y + Inches(1.4), 
                              Inches(5.4), Inches(0.5),
                              desc, font_size=14, color=TEXT_MUTED)
        
        self.slide_count += 1
        return slide
    
    def add_chapter(self, number: str, title: str, subtitle: str,
                   accent: RGBColor = NEON_PURPLE):
        """添加章节过渡页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_bg(slide)
        
        # Large background number (watermark style)
        self._add_text_box(slide, Inches(0), Inches(1.5), Inches(13.333), Inches(2.5),
                          number, font_size=180, bold=True, 
                          color=RGBColor(0x30, 0x30, 0x30), 
                          align=PP_ALIGN.CENTER)
        
        # Title
        self._add_text_box(slide, Inches(0), Inches(3.0), Inches(13.333), Inches(1.0),
                          title, font_size=56, bold=True, 
                          color=TEXT_PRIMARY, align=PP_ALIGN.CENTER)
        
        # Accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                      Inches(5.8), Inches(4.0), 
                                      Inches(1.733), Pt(4))
        line.fill.solid()
        line.fill.fore_color.rgb = accent
        line.line.fill.background()
        
        # Subtitle
        self._add_text_box(slide, Inches(0), Inches(4.3), Inches(13.333), Inches(0.6),
                          subtitle, font_size=20, 
                          color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        
        self.slide_count += 1
        return slide
    
    def add_two_col(self, title: str, 
                   left_title: str, left_items: List[str],
                   right_title: str, right_items: List[str],
                   subtitle: Optional[str] = None,
                   left_accent: RGBColor = NEON_PURPLE,
                   right_accent: RGBColor = NEON_EMERALD,
                   quote: Optional[str] = None):
        """添加双栏内容页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_bg(slide)
        
        # Title
        self._add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                          title, font_size=32, bold=True, color=TEXT_PRIMARY)
        
        # Title underline
        underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                          Inches(0.5), Inches(0.95), 
                                          Inches(2), Pt(2))
        underline.fill.solid()
        underline.fill.fore_color.rgb = NEON_PURPLE
        underline.line.fill.background()
        
        y_offset = Inches(1.2)
        if subtitle:
            self._add_text_box(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.4),
                              subtitle, font_size=14, color=TEXT_MUTED)
            y_offset = Inches(1.5)
        
        # Left column
        card = self._add_rounded_rect(slide, Inches(0.5), y_offset, 
                                     Inches(5.8), Inches(4.5), 
                                     CARD_BG, left_accent, Pt(1))
        self._add_text_box(slide, Inches(0.7), y_offset + Inches(0.15), 
                          Inches(5.4), Inches(0.5),
                          left_title, font_size=20, bold=True, 
                          color=left_accent)
        self._add_bullet_text(slide, Inches(0.7), y_offset + Inches(0.7), 
                             Inches(5.4), Inches(3.5),
                             left_items, font_size=14, 
                             color=TEXT_SECONDARY, bullet_color=left_accent)
        
        # Right column
        card = self._add_rounded_rect(slide, Inches(6.8), y_offset, 
                                     Inches(5.8), Inches(4.5), 
                                     CARD_BG, right_accent, Pt(1))
        self._add_text_box(slide, Inches(7.0), y_offset + Inches(0.15), 
                          Inches(5.4), Inches(0.5),
                          right_title, font_size=20, bold=True, 
                          color=right_accent)
        self._add_bullet_text(slide, Inches(7.0), y_offset + Inches(0.7), 
                             Inches(5.4), Inches(3.5),
                             right_items, font_size=14, 
                             color=TEXT_SECONDARY, bullet_color=right_accent)
        
        # Quote
        if quote:
            quote_y = Inches(6.2)
            card = self._add_rounded_rect(slide, Inches(0.5), quote_y, 
                                         Inches(12.1), Inches(0.9), 
                                         RGBColor(0x2D, 0x1B, 0x4E), 
                                         NEON_PURPLE, Pt(0.5))
            self._add_text_box(slide, Inches(0.7), quote_y + Inches(0.15), 
                              Inches(11.7), Inches(0.6),
                              f'"{quote}"', font_size=13, 
                              color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
        
        self.slide_count += 1
        return slide
    
    def add_list(self, title: str, items: List[str],
                subtitle: Optional[str] = None,
                accent: RGBColor = NEON_PURPLE):
        """添加列表内容页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_bg(slide)
        
        # Title
        self._add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                          title, font_size=32, bold=True, color=TEXT_PRIMARY)
        
        # Title underline
        underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                          Inches(0.5), Inches(0.95), 
                                          Inches(2), Pt(2))
        underline.fill.solid()
        underline.fill.fore_color.rgb = accent
        underline.line.fill.background()
        
        y_offset = Inches(1.2)
        if subtitle:
            self._add_text_box(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.4),
                              subtitle, font_size=14, color=TEXT_MUTED)
            y_offset = Inches(1.5)
        
        # Content card
        card = self._add_rounded_rect(slide, Inches(0.5), y_offset, 
                                     Inches(12.1), Inches(4.5), 
                                     CARD_BG, accent, Pt(1))
        self._add_bullet_text(slide, Inches(0.7), y_offset + Inches(0.3), 
                             Inches(11.7), Inches(4.0),
                             items, font_size=16, 
                             color=TEXT_SECONDARY, bullet_color=accent)
        
        self.slide_count += 1
        return slide
    
    def add_code(self, title: str, code_text: str, highlight: str,
                subtitle: Optional[str] = None,
                accent: RGBColor = NEON_PURPLE):
        """添加代码展示页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_bg(slide)
        
        # Title
        self._add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                          title, font_size=32, bold=True, color=TEXT_PRIMARY)
        
        # Title underline
        underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                          Inches(0.5), Inches(0.95), 
                                          Inches(2), Pt(2))
        underline.fill.solid()
        underline.fill.fore_color.rgb = accent
        underline.line.fill.background()
        
        y_offset = Inches(1.2)
        if subtitle:
            self._add_text_box(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.4),
                              subtitle, font_size=14, color=TEXT_MUTED)
            y_offset = Inches(1.5)
        
        # Code card
        code_height = Inches(4.2)
        card = self._add_rounded_rect(slide, Inches(0.5), y_offset, 
                                     Inches(12.1), code_height, 
                                     RGBColor(0x0D, 0x0D, 0x10), 
                                     accent, Pt(0.5))
        
        # Code text
        self._add_text_box(slide, Inches(0.7), y_offset + Inches(0.15), 
                          Inches(11.7), code_height - Inches(0.3),
                          code_text, font_size=10, 
                          color=TEXT_SECONDARY, font_name="Consolas")
        
        # Highlight text
        highlight_y = y_offset + code_height + Inches(0.15)
        self._add_text_box(slide, Inches(0.5), highlight_y, Inches(12), Inches(0.4),
                          f"亮点：{highlight}", font_size=12, 
                          color=accent)
        
        self.slide_count += 1
        return slide
    
    def add_grid(self, title: str, items: List[Tuple[str, str, RGBColor]],
                subtitle: Optional[str] = None,
                cols: int = 4):
        """添加网格卡片页
        
        Args:
            items: List of (title, description, accent_color)
            cols: Number of columns (4 or 5)
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_bg(slide)
        
        # Title
        self._add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                          title, font_size=32, bold=True, color=TEXT_PRIMARY)
        
        # Title underline
        underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                          Inches(0.5), Inches(0.95), 
                                          Inches(2), Pt(2))
        underline.fill.solid()
        underline.fill.fore_color.rgb = NEON_PURPLE
        underline.line.fill.background()
        
        y_offset = Inches(1.2)
        if subtitle:
            self._add_text_box(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.4),
                              subtitle, font_size=14, color=TEXT_MUTED)
            y_offset = Inches(1.5)
        
        if cols == 5:
            card_width = Inches(2.2)
            card_height = Inches(2.8)
            gap = Inches(0.25)
            start_x = Inches(0.5)
            
            for i, (item_title, desc, accent) in enumerate(items):
                x = start_x + i * (card_width + gap)
                card = self._add_rounded_rect(slide, x, y_offset, 
                                             card_width, card_height, 
                                             CARD_BG, accent, Pt(1))
                self._add_text_box(slide, x + Inches(0.1), y_offset + Inches(0.15), 
                                  Inches(2.0), Inches(0.5),
                                  item_title, font_size=18, bold=True, 
                                  color=accent, align=PP_ALIGN.CENTER)
                self._add_text_box(slide, x + Inches(0.1), y_offset + Inches(0.7), 
                                  Inches(2.0), Inches(1.8),
                                  desc, font_size=12, 
                                  color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        else:  # 4 cols (2x2 grid)
            card_width = Inches(5.8)
            card_height = Inches(2.2)
            gap_x = Inches(0.4)
            gap_y = Inches(0.3)
            start_x = Inches(0.5)
            
            for i, (item_title, desc, accent) in enumerate(items):
                row = i // 2
                col = i % 2
                x = start_x + col * (card_width + gap_x)
                y = y_offset + row * (card_height + gap_y)
                card = self._add_rounded_rect(slide, x, y, 
                                             card_width, card_height, 
                                             CARD_BG, accent, Pt(1))
                self._add_text_box(slide, x + Inches(0.2), y + Inches(0.15), 
                                  Inches(5.4), Inches(0.4),
                                  item_title, font_size=20, bold=True, 
                                  color=accent)
                self._add_text_box(slide, x + Inches(0.2), y + Inches(0.6), 
                                  Inches(5.4), Inches(1.3),
                                  desc, font_size=13, 
                                  color=TEXT_SECONDARY)
        
        self.slide_count += 1
        return slide
    
    def add_three_col(self, title: str, 
                     cols: List[Dict[str, Any]],
                     subtitle: Optional[str] = None):
        """添加三栏内容页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_bg(slide)
        
        # Title
        self._add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                          title, font_size=32, bold=True, color=TEXT_PRIMARY)
        
        # Title underline
        underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                          Inches(0.5), Inches(0.95), 
                                          Inches(2), Pt(2))
        underline.fill.solid()
        underline.fill.fore_color.rgb = NEON_PURPLE
        underline.line.fill.background()
        
        y_offset = Inches(1.2)
        if subtitle:
            self._add_text_box(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.4),
                              subtitle, font_size=14, color=TEXT_MUTED)
            y_offset = Inches(1.5)
        
        col_width = Inches(3.8)
        gap = Inches(0.3)
        start_x = Inches(0.5)
        
        for i, col_data in enumerate(cols):
            x = start_x + i * (col_width + gap)
            accent = col_data.get("accent", NEON_PURPLE)
            card = self._add_rounded_rect(slide, x, y_offset, 
                                         col_width, Inches(4.5), 
                                         CARD_BG, accent, Pt(1))
            self._add_text_box(slide, x + Inches(0.15), y_offset + Inches(0.15), 
                              Inches(3.5), Inches(0.5),
                              col_data["title"], font_size=20, bold=True, 
                              color=accent)
            self._add_text_box(slide, x + Inches(0.15), y_offset + Inches(0.7), 
                              Inches(3.5), Inches(1.0),
                              col_data["content"], font_size=13, 
                              color=TEXT_SECONDARY)
            self._add_text_box(slide, x + Inches(0.15), y_offset + Inches(1.8), 
                              Inches(3.5), Inches(2.0),
                              col_data.get("example", ""), font_size=11, 
                              color=TEXT_MUTED)
        
        self.slide_count += 1
        return slide
    
    def add_stats(self, title: str, 
                 left_data: Dict, right_data: Dict,
                 stats: Optional[List[Tuple[str, str]]] = None,
                 big_stat: Optional[Tuple[str, str]] = None,
                 subtitle: Optional[str] = None):
        """添加数据统计页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_bg(slide)
        
        # Title
        self._add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                          title, font_size=32, bold=True, color=TEXT_PRIMARY)
        
        # Title underline
        underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                          Inches(0.5), Inches(0.95), 
                                          Inches(2), Pt(2))
        underline.fill.solid()
        underline.fill.fore_color.rgb = NEON_BLUE
        underline.line.fill.background()
        
        y_offset = Inches(1.2)
        if subtitle:
            self._add_text_box(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.4),
                              subtitle, font_size=14, color=TEXT_MUTED)
            y_offset = Inches(1.5)
        
        # Left column
        if left_data:
            card = self._add_rounded_rect(slide, Inches(0.5), y_offset, 
                                         Inches(5.8), Inches(4.5), 
                                         CARD_BG, NEON_BLUE, Pt(1))
            self._add_text_box(slide, Inches(0.7), y_offset + Inches(0.15), 
                              Inches(5.4), Inches(0.5),
                              left_data["title"], font_size=20, bold=True, 
                              color=NEON_BLUE)
            self._add_bullet_text(slide, Inches(0.7), y_offset + Inches(0.7), 
                                 Inches(5.4), Inches(3.5),
                                 left_data["items"], font_size=14, 
                                 color=TEXT_SECONDARY, bullet_color=NEON_BLUE)
        
        # Right column
        if right_data:
            card = self._add_rounded_rect(slide, Inches(6.8), y_offset, 
                                         Inches(5.8), Inches(4.5), 
                                         CARD_BG, NEON_EMERALD, Pt(1))
            self._add_text_box(slide, Inches(7.0), y_offset + Inches(0.15), 
                              Inches(5.4), Inches(0.5),
                              right_data["title"], font_size=20, bold=True, 
                              color=NEON_EMERALD)
            self._add_bullet_text(slide, Inches(7.0), y_offset + Inches(0.7), 
                                 Inches(5.4), Inches(3.5),
                                 right_data["items"], font_size=14, 
                                 color=TEXT_SECONDARY, bullet_color=NEON_EMERALD)
        
        # Stats cards
        if stats:
            stats_y = Inches(6.0)
            for i, (value, label) in enumerate(stats):
                x = Inches(3.5 + i * 3.5)
                card = self._add_rounded_rect(slide, x, stats_y, 
                                             Inches(3.0), Inches(1.0), 
                                             RGBColor(0x0D, 0x0D, 0x10), 
                                             NEON_BLUE, Pt(0.5))
                self._add_text_box(slide, x, stats_y + Inches(0.1), 
                                  Inches(3.0), Inches(0.5),
                                  value, font_size=28, bold=True, 
                                  color=NEON_BLUE, align=PP_ALIGN.CENTER)
                self._add_text_box(slide, x, stats_y + Inches(0.5), 
                                  Inches(3.0), Inches(0.4),
                                  label, font_size=12, 
                                  color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        
        # Big stat
        if big_stat:
            value, label = big_stat
            self._add_text_box(slide, Inches(0), Inches(6.0), Inches(13.333), Inches(0.8),
                              "RAG 模式显著降低幻觉，", font_size=20, 
                              color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
            self._add_text_box(slide, Inches(0), Inches(6.6), Inches(13.333), Inches(0.8),
                              f"忠实度 {value}", font_size=36, bold=True, 
                              color=NEON_EMERALD, align=PP_ALIGN.CENTER)
        
        self.slide_count += 1
        return slide
    
    def save(self, output_path: str):
        """保存PPT文件"""
        self.prs.save(output_path)
        print(f"PPT saved: {output_path}")
        print(f"Total slides: {self.slide_count}")
        return output_path
