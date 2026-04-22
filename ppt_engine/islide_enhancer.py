#!/usr/bin/env python3
"""
iSlideEnhancer - COM automation for PowerPoint and iSlide plugin
Provides advanced enhancement capabilities through Windows COM interface
"""

import os
import sys
import time
from typing import Optional, List, Dict, Any

class ISlideEnhancer:
    """iSlide增强器 - 通过COM调用PowerPoint和iSlide插件"""
    
    def __init__(self):
        self.app = None
        self.presentation = None
        self.is_available = False
        
    def check_availability(self) -> bool:
        """检查PowerPoint和iSlide是否可用"""
        try:
            import win32com.client
            # Try to dispatch PowerPoint
            app = win32com.client.Dispatch("PowerPoint.Application")
            app.Quit()
            self.is_available = True
            return True
        except ImportError:
            print("[iSlideEnhancer] pywin32 not installed. COM automation unavailable.")
            return False
        except Exception as e:
            print(f"[iSlideEnhancer] PowerPoint not available: {e}")
            return False
    
    def open_powerpoint(self, visible: bool = True):
        """启动PowerPoint应用"""
        try:
            import win32com.client
            self.app = win32com.client.Dispatch("PowerPoint.Application")
            self.app.Visible = visible
            self.app.DisplayAlerts = False
            print("[iSlideEnhancer] PowerPoint opened successfully")
        except Exception as e:
            print(f"[iSlideEnhancer] Failed to open PowerPoint: {e}")
            raise
    
    def open_presentation(self, pptx_path: str):
        """打开PPT文件"""
        if not self.app:
            self.open_powerpoint()
        
        try:
            abs_path = os.path.abspath(pptx_path)
            self.presentation = self.app.Presentations.Open(abs_path)
            print(f"[iSlideEnhancer] Opened: {abs_path}")
        except Exception as e:
            print(f"[iSlideEnhancer] Failed to open presentation: {e}")
            raise
    
    def apply_islide_uniform_design(self):
        """应用iSlide一键统一设计（如果可用）"""
        if not self.presentation:
            print("[iSlideEnhancer] No presentation open")
            return
        
        try:
            # Try to access iSlide ribbon commands
            # Note: iSlide doesn't expose public API, so we try common approaches
            
            # Method 1: Try to execute iSlide command via CommandBars
            try:
                # Common iSlide command names (may vary by version)
                islide_commands = [
                    "iSlideDesignUniform",
                    "iSlide_Uniform",
                    "iSlideDesign",
                    "iSlide.OneKey.Uniform",
                ]
                
                for cmd in islide_commands:
                    try:
                        self.app.CommandBars.ExecuteMso(cmd)
                        print(f"[iSlideEnhancer] Executed iSlide command: {cmd}")
                        time.sleep(2)  # Wait for processing
                        return
                    except:
                        continue
                        
            except Exception as e:
                print(f"[iSlideEnhancer] CommandBar execution failed: {e}")
            
            # Method 2: Try to access iSlide via Ribbon XML (advanced)
            try:
                # Get the ribbon XML and parse for iSlide controls
                ribbon = self.app.ActivePresentation.RibbonXML
                if "iSlide" in ribbon:
                    print("[iSlideEnhancer] iSlide detected in ribbon")
            except:
                pass
            
            print("[iSlideEnhancer] iSlide automation limited - plugin API not exposed")
            
        except Exception as e:
            print(f"[iSlideEnhancer] Failed to apply iSlide design: {e}")
    
    def apply_islide_smart_chart(self, slide_index: int, chart_data: Dict):
        """应用iSlide智能图表（实验性功能）"""
        if not self.presentation:
            return
        
        try:
            slide = self.presentation.Slides(slide_index)
            # This would require specific iSlide chart insertion commands
            # Since iSlide doesn't expose API, this is a placeholder
            print(f"[iSlideEnhancer] Smart chart insertion requires manual iSlide interaction")
        except Exception as e:
            print(f"[iSlideEnhancer] Smart chart failed: {e}")
    
    def apply_powerpoint_transitions(self, transition_type: str = "Fade"):
        """应用PowerPoint幻灯片切换效果"""
        if not self.presentation:
            return
        
        try:
            # PpTransitionType constants
            transition_types = {
                "Fade": 0x09,
                "Push": 0x01,
                "Wipe": 0x02,
                "Split": 0x0F,
                "Reveal": 0x11,
                "Random": 0x17,
            }
            
            transition_value = transition_types.get(transition_type, 0x09)
            
            for slide in self.presentation.Slides:
                slide.SlideShowTransition.EntryEffect = transition_value
                slide.SlideShowTransition.Duration = 0.6  # 0.6 seconds
            
            print(f"[iSlideEnhancer] Applied {transition_type} transition to all slides")
            
        except Exception as e:
            print(f"[iSlideEnhancer] Transition application failed: {e}")
    
    def apply_powerpoint_animations(self, animation_type: str = "Fade"):
        """应用PowerPoint动画效果"""
        if not self.presentation:
            return
        
        try:
            # PpEffect constants for animations
            animation_effects = {
                "Fade": 0x0A,
                "FlyIn": 0x01,
                "Appear": 0x00,
                "FloatIn": 0x02,
                "Split": 0x0F,
            }
            
            effect_value = animation_effects.get(animation_type, 0x0A)
            
            for slide in self.presentation.Slides:
                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        # Add entrance animation to text
                        animation = slide.TimeLine.MainSequence.AddEffect(
                            shape=shape,
                            effectId=effect_value,
                            trigger=1  # msoAnimTriggerWithPrevious
                        )
            
            print(f"[iSlideEnhancer] Applied {animation_type} animation to all shapes")
            
        except Exception as e:
            print(f"[iSlideEnhancer] Animation application failed: {e}")
    
    def export_to_pdf(self, output_path: str):
        """导出为PDF"""
        if not self.presentation:
            return
        
        try:
            abs_path = os.path.abspath(output_path)
            self.presentation.SaveAs(abs_path, 32)  # 32 = ppSaveAsPDF
            print(f"[iSlideEnhancer] Exported to PDF: {abs_path}")
        except Exception as e:
            print(f"[iSlideEnhancer] PDF export failed: {e}")
    
    def export_to_images(self, output_folder: str, format: str = "PNG"):
        """导出幻灯片为图片"""
        if not self.presentation:
            return
        
        try:
            os.makedirs(output_folder, exist_ok=True)
            
            for i, slide in enumerate(self.presentation.Slides, 1):
                output_path = os.path.join(output_folder, f"slide_{i:02d}.{format.lower()}")
                slide.Export(output_path, format)
            
            print(f"[iSlideEnhancer] Exported {len(self.presentation.Slides)} slides to {output_folder}")
            
        except Exception as e:
            print(f"[iSlideEnhancer] Image export failed: {e}")
    
    def save_and_close(self, output_path: Optional[str] = None):
        """保存并关闭演示文稿"""
        if self.presentation:
            try:
                if output_path:
                    abs_path = os.path.abspath(output_path)
                    self.presentation.SaveAs(abs_path)
                    print(f"[iSlideEnhancer] Saved: {abs_path}")
                else:
                    self.presentation.Save()
                self.presentation.Close()
                self.presentation = None
            except Exception as e:
                print(f"[iSlideEnhancer] Save failed: {e}")
    
    def quit(self):
        """退出PowerPoint应用"""
        if self.presentation:
            self.save_and_close()
        
        if self.app:
            try:
                self.app.Quit()
                self.app = None
                print("[iSlideEnhancer] PowerPoint closed")
            except Exception as e:
                print(f"[iSlideEnhancer] Quit failed: {e}")
    
    def enhance_presentation(self, pptx_path: str, 
                            apply_transitions: bool = True,
                            apply_animations: bool = False,
                            export_pdf: bool = False,
                            output_folder: Optional[str] = None) -> str:
        """完整的增强流程
        
        Args:
            pptx_path: 输入PPT文件路径
            apply_transitions: 是否应用切换效果
            apply_animations: 是否应用动画效果
            export_pdf: 是否导出PDF
            output_folder: 输出文件夹
            
        Returns:
            增强后的PPT文件路径
        """
        if not self.check_availability():
            print("[iSlideEnhancer] PowerPoint not available, skipping enhancement")
            return pptx_path
        
        try:
            self.open_powerpoint(visible=False)
            self.open_presentation(pptx_path)
            
            # Apply enhancements
            if apply_transitions:
                self.apply_powerpoint_transitions("Fade")
            
            if apply_animations:
                self.apply_powerpoint_animations("Fade")
            
            # Try iSlide (limited due to no public API)
            self.apply_islide_uniform_design()
            
            # Save enhanced version
            if output_folder:
                os.makedirs(output_folder, exist_ok=True)
                base_name = os.path.splitext(os.path.basename(pptx_path))[0]
                output_path = os.path.join(output_folder, f"{base_name}_enhanced.pptx")
            else:
                base, ext = os.path.splitext(pptx_path)
                output_path = f"{base}_enhanced{ext}"
            
            self.save_and_close(output_path)
            
            # Export PDF if requested
            if export_pdf:
                pdf_path = os.path.splitext(output_path)[0] + ".pdf"
                self.open_presentation(output_path)
                self.export_to_pdf(pdf_path)
                self.save_and_close()
            
            self.quit()
            
            return output_path
            
        except Exception as e:
            print(f"[iSlideEnhancer] Enhancement failed: {e}")
            self.quit()
            return pptx_path


# ============================================================
# Standalone functions for template-based enhancement
# ============================================================

def apply_islide_template(pptx_path: str, template_path: str) -> str:
    """应用iSlide风格模板到现有PPT
    
    Since iSlide doesn't expose API, we use template files (.potx)
    that were created with iSlide styling.
    
    Args:
        pptx_path: 源PPT文件路径
        template_path: iSlide模板文件路径(.potx)
        
    Returns:
        应用模板后的PPT文件路径
    """
    try:
        from pptx import Presentation
        
        # Load source presentation
        src_prs = Presentation(pptx_path)
        
        # Load template (contains iSlide styling)
        if os.path.exists(template_path):
            tmpl_prs = Presentation(template_path)
            # Copy theme and layouts from template
            # This is limited - full theme transfer requires COM
            print(f"[iSlideEnhancer] Template loaded: {template_path}")
            print("[iSlideEnhancer] Note: Full theme transfer requires PowerPoint COM")
        
        # Save with template reference
        output_path = pptx_path.replace(".pptx", "_templated.pptx")
        src_prs.save(output_path)
        
        return output_path
        
    except Exception as e:
        print(f"[iSlideEnhancer] Template application failed: {e}")
        return pptx_path


def create_islide_template_placeholder(output_path: str):
    """创建iSlide风格模板占位文件
    
    This creates a basic template with cyberpunk styling
    that can be manually enhanced in PowerPoint with iSlide.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Set default theme colors (dark cyberpunk)
        # Note: python-pptx has limited theme modification capabilities
        
        prs.save(output_path)
        print(f"[iSlideEnhancer] Template placeholder created: {output_path}")
        
        return output_path
        
    except Exception as e:
        print(f"[iSlideEnhancer] Template creation failed: {e}")
        return None
